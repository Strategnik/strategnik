#!/usr/bin/env python3
"""
PLG Framework — Infographic Deck Generator
Generates one full-bleed infographic image per slide via fal.ai (nano-banana-2),
then assembles into a PPTX.
"""

import json
import os
import sys
import time
import requests
from pathlib import Path
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Emu

# --- Config ---
FAL_KEY = os.environ.get("FAL_KEY")
if not FAL_KEY:
    sys.exit("FAL_KEY not set. Export it first.")

MODEL = "fal-ai/nano-banana-2"
FAL_URL = f"https://fal.run/{MODEL}"
OUTPUT_DIR = Path(__file__).parent / "slide-images"
OUTPUT_DIR.mkdir(exist_ok=True)

# 16:9 at high res
IMG_W, IMG_H = 1792, 1024

# --- Style Directive ---
STYLE = """STYLE — HBR / MIT Sloan editorial infographic:
- Background: #F4EFE6 warm cream with very subtle paper texture
- Primary text: #1A1A1A soft black
- Headlines: #1B3A5B deep navy, transitional serif font (like Tiempos or Georgia), title case
- Accent numerals/callouts: #C9402F editorial red-orange, used sparingly
- Secondary text: #8A8680 warm gray
- Dividers/fills: #D9D4CB light warm gray
- Body text: clean geometric sans-serif (like Inter or Helvetica)
- Data labels: monospace font, lowercase for variable names
- Illustration: flat editorial line art, 1-1.5pt strokes, no gradients, no shadows, no 3D, no glossy effects
- Layout: generous whitespace, clear hierarchy, print-editorial feel
- Do NOT include any slide numbers or page numbers in the image
- This should look like it belongs in a printed strategy journal"""

TEXT_RULES = """TEXT RENDERING (mandatory):
Render ALL text listed below VERBATIM. Do not add, omit, paraphrase, or misspell any text.
Do not add logos, URLs, watermarks, or decorative text not specified.
If text won't fit legibly, reduce visual complexity — never alter the text."""

# --- Slide Definitions ---
SLIDES = [
    {
        "id": "slide-01-title",
        "prompt": f"""{STYLE}

Create a title slide for a strategy presentation. Full-bleed editorial layout.

COMPOSITION:
- Large serif headline dominating the left two-thirds of the canvas
- A thin red-orange horizontal rule (3px) above or below the headline
- Small uppercase monospace eyebrow text above the headline
- Italic serif subtitle below the headline
- Author/brand line at bottom in small sans-serif
- Right third: subtle decorative element — fine diagonal hatching lines in light gray, or thin geometric line pattern
- Generous negative space throughout

TEXT TO RENDER:
Eyebrow: "A PLG OPERATING FRAMEWORK"
Headline: "The PLG Framework for Developer Products"
Subtitle: "Three systems that convert free users into qualified pipeline."
Byline: "Nick Talbert — Strategnik"

{TEXT_RULES}"""
    },
    {
        "id": "slide-02-problem",
        "prompt": f"""{STYLE}

Create an editorial infographic slide about a business problem.

COMPOSITION:
- Top-left: small monospace label "THE PROBLEM" and a short navy rule beneath it
- Large serif headline: the main message
- Below: a paragraph of supporting text in gray sans-serif
- A pull quote in italic serif with a red-orange left border bar (3px vertical line)
- Clean, text-forward layout with generous whitespace — this is a narrative slide, not a data slide

TEXT TO RENDER:
Label: "THE PROBLEM"
Headline: "When Frictionless Breaks"
Body: "The PLG orthodoxy says remove all friction. But when your API is too easy to hit, you get usage data that looks great and predicts nothing."
Pull quote: "The signal-to-noise ratio had collapsed, and nobody could tell the difference."

{TEXT_RULES}"""
    },
    {
        "id": "slide-03-cascading-damage",
        "prompt": f"""{STYLE}

Create an infographic slide showing three cascading consequences. Editorial layout.

COMPOSITION:
- Top: monospace label "WHAT HAPPENS WHEN SIGNAL COLLAPSES" with a short navy rule
- Three stacked rows, each separated by a thin warm-gray horizontal rule
- Each row has a bold serif sub-headline on the left and a supporting description in gray sans-serif on the right or below
- No icons or illustrations — this is a text-driven editorial layout
- Clear visual hierarchy between the three items

TEXT TO RENDER:
Label: "WHAT HAPPENS WHEN SIGNAL COLLAPSES"
Row 1 headline: "Sales wastes cycles on phantom intent"
Row 1 body: "PQLs stop predicting conversion. Sales stops trusting the system. Real prospects get buried."
Row 2 headline: "Usage metrics become vanity metrics"
Row 2 body: "MAU goes up. Revenue doesn't. The metric becomes self-reinforcing and self-deceiving."
Row 3 headline: "ICP definition blurs"
Row 3 body: "Product decisions get made on aggregate behavior of a user base that's 80% non-buyers."

{TEXT_RULES}"""
    },
    {
        "id": "slide-04-thesis",
        "prompt": f"""{STYLE}

EXCEPTION: This slide uses a DARK BACKGROUND. Background: #1B3A5B deep navy. All text is cream #F4EFE6 or white. The red-orange #C9402F accent is used for the rule.

Create a bold thesis statement slide. Dark navy background.

COMPOSITION:
- Centered or left-aligned layout
- Small monospace label "THE THESIS" in muted cream
- A short red-orange rule beneath the label
- Very large serif headline — the core thesis statement
- Below: a supporting paragraph in lighter cream
- Below that: a thin rule, then a secondary quote in slightly smaller text
- This should feel like a keynote moment — bold, confident, dramatic

TEXT TO RENDER:
Label: "THE THESIS"
Headline: "Friction as Signal Architecture"
Supporting text: "Make it easy to use your product for the right reasons, and slightly less easy for the wrong ones."
Secondary text: "A twelve-field form and a work email requirement are both 'friction.' One is process friction imposed on the buyer. The other is signal friction that benefits everyone."

{TEXT_RULES}"""
    },
    {
        "id": "slide-05-three-systems",
        "prompt": f"""{STYLE}

Create an infographic slide showing a three-part framework overview.

COMPOSITION:
- Top: monospace label "THE FRAMEWORK" with a short navy rule
- Headline: "Three Systems"
- Below: three equal columns or three large blocks arranged horizontally
- Each block has: a large red-orange numeral (01, 02, 03), a navy serif title, and a short description in gray sans-serif
- Thin connecting arrows or lines between the three blocks flowing left to right
- Each block could have a small abstract icon: Block 01 = filter/funnel shape, Block 02 = gauge/scale shape, Block 03 = radial burst shape
- Clean grid layout with balanced spacing

TEXT TO RENDER:
Label: "THE FRAMEWORK"
Headline: "Three Systems"
Block 01 numeral: "01"
Block 01 title: "Signal Architecture"
Block 01 description: "Friction points that capture buying intent data"
Block 02 numeral: "02"
Block 02 title: "PQL Scoring"
Block 02 description: "Variables that identify when a user is ready for sales"
Block 03 numeral: "03"
Block 03 title: "Brand Infrastructure"
Block 03 description: "GitHub repos and technical content that multiply touchpoints"

{TEXT_RULES}"""
    },
    {
        "id": "slide-06-signup-gate",
        "prompt": f"""{STYLE}

Create an infographic slide about the first signal architecture friction point.

COMPOSITION:
- Two-column layout: text on left, visual on right
- Left column: small monospace "SIGNAL ARCHITECTURE" label, large red-orange numeral "01", navy serif title, body text in sans-serif, and two monospace variable tags in rounded gray pill shapes at bottom
- Right column: an abstract editorial diagram showing a filter concept — a row of circles meeting a horizontal line, some passing through (allowed), some blocked with an X mark. Simple, diagrammatic, flat line art.
- Below the diagram: a pull quote in italic serif with a red-orange left border

TEXT TO RENDER:
Label: "SIGNAL ARCHITECTURE"
Numeral: "01"
Title: "Signup Gate"
Body: "Work email + company + stated use case. Filters 40% of non-buyers before they cost you anything downstream."
Pull quote: "The developers who are serious don't mind. The ones who bounce were never going to convert."
Variable tag 1: "has_work_email"
Variable tag 2: "stated_use_case"

{TEXT_RULES}"""
    },
    {
        "id": "slide-07-integration-boundary",
        "prompt": f"""{STYLE}

Create an infographic slide about the second signal architecture friction point.

COMPOSITION:
- Two-column layout: text on left, visual diagram on right
- Left column: small monospace "SIGNAL ARCHITECTURE" label, large red-orange numeral "02", navy serif title, body text, and two monospace variable tags in gray pill shapes
- Right column: a diagram showing two adjacent rectangular zones separated by a bold vertical dashed line. Left zone filled with light gray #D9D4CB and labeled "SANDBOX". Right zone has a thin navy border on white/cream and labeled "PRODUCTION". A small navy circle with a curved arrow crosses from left zone to right zone.
- This represents the boundary between exploration and real commitment

TEXT TO RENDER:
Label: "SIGNAL ARCHITECTURE"
Numeral: "02"
Title: "Integration Boundary"
Body: "Sandbox is frictionless. Production requires a verification step. Where casual exploration becomes real commitment."
Zone left label: "SANDBOX"
Zone right label: "PRODUCTION"
Variable tag 1: "production_integration_date"
Variable tag 2: "integration_depth_score"

{TEXT_RULES}"""
    },
    {
        "id": "slide-08-conversion-wall",
        "prompt": f"""{STYLE}

Create an infographic slide about the third signal architecture friction point.

COMPOSITION:
- Two-column layout: text on left, visual on right
- Left column: small monospace "SIGNAL ARCHITECTURE" label, large red-orange numeral "03", navy serif title, body text, and two monospace variable tags in gray pill shapes
- Right column: a minimal line chart diagram showing an ascending curve that hits a horizontal ceiling line (the wall). A thin dashed horizontal rule sits at about 70% height, labeled "70% free-tier limit" in monospace. The curve should be drawn in navy, the ceiling in darker weight.
- Pull quote below the chart area

TEXT TO RENDER:
Label: "SIGNAL ARCHITECTURE"
Numeral: "03"
Title: "Conversion Wall"
Body: "Usage limits plus feature gates. Tuned to convert growth trajectory, not cap adoption."
Chart label: "70% free-tier limit"
Pull quote: "If more than 30% of active users never hit the wall, it's set too high."
Variable tag 1: "usage_limit_proximity"
Variable tag 2: "pricing_page_visits"

{TEXT_RULES}"""
    },
    {
        "id": "slide-09-team-expansion",
        "prompt": f"""{STYLE}

Create an infographic slide about the fourth signal architecture friction point.

COMPOSITION:
- Two-column layout: text on left, visual on right
- Left column: small monospace "SIGNAL ARCHITECTURE" label, large red-orange numeral "04", navy serif title, body text, two monospace variable tags in gray pill shapes stacked vertically
- Right column: a network node diagram. ONE large navy circle on the far left representing a solo developer. FIVE navy circles clustered on the right in a loose group, connected by thin lines to the single circle (hub-and-spoke pattern). This shows one dev expanding to a team.
- Pull quote below

TEXT TO RENDER:
Label: "SIGNAL ARCHITECTURE"
Numeral: "04"
Title: "Team Expansion Gate"
Body: "Collaboration behind paid tier. Invite limits on free. One dev is an experiment. Five devs is an evaluation."
Pull quote: "A team of five on a free plan is either a serious evaluation or a company working around your paywall."
Variable tag 1: "team_size"
Variable tag 2: "admin_engagement_score"

{TEXT_RULES}"""
    },
    {
        "id": "slide-10-pql-intro",
        "prompt": f"""{STYLE}

EXCEPTION: This slide uses a DARK BACKGROUND. Background: #1B3A5B deep navy. All text is cream #F4EFE6. The red-orange #C9402F accent is used for the rule.

Create a section divider slide for the PQL scoring system. Dark navy background.

COMPOSITION:
- Monospace label "SYSTEM 02 — PQL SCORING" in muted cream
- A short red-orange rule beneath
- Large serif headline
- Italic subtitle beneath in lighter cream
- A very large, very faint numeral "02" in the background (about 10% opacity) as a watermark element
- Abstract diagram on the right: three horizontal lines of different thicknesses entering from the left, converging into a single circular node on the right (representing weighted inputs producing a score)

TEXT TO RENDER:
Label: "SYSTEM 02 — PQL SCORING"
Headline: "Scoring for Production Dependency"
Subtitle: "The question every variable answers: is this user showing signs of production dependency plus enterprise need?"

{TEXT_RULES}"""
    },
    {
        "id": "slide-11-variables",
        "prompt": f"""{STYLE}

Create an infographic slide showing three categories of PQL variables in a three-column layout.

COMPOSITION:
- Top: monospace label "PQL VARIABLES" with a short navy rule
- Three equal columns separated by thin vertical rules (1px warm gray lines)
- Each column has: a monospace uppercase category header in gray, a LARGE serif percentage numeral in navy below, and a bulleted list of variable names in monospace below that
- The percentage numerals (40%, 40%, 20%) should be visually prominent — these are the weights
- Clean, data-forward layout

TEXT TO RENDER:
Label: "PQL VARIABLES"
Column 1 header: "PRODUCTION USAGE"
Column 1 weight: "40%"
Column 1 variables: "api_calls_last_30d" "error_rate" "integration_points" "usage_consistency"
Column 2 header: "ENTERPRISE INTENT"
Column 2 weight: "40%"
Column 2 variables: "team_size" "sso_attempt" "usage_limit_proximity" "domain_firmographic"
Column 3 header: "ENGAGEMENT"
Column 3 weight: "20%"
Column 3 variables: "docs_depth" "community_participation" "support_ticket_quality" "pricing_page_visits"

{TEXT_RULES}"""
    },
    {
        "id": "slide-12-formula",
        "prompt": f"""{STYLE}

Create an infographic slide displaying a scoring formula as a typographic centerpiece.

COMPOSITION:
- Centered layout with generous whitespace
- Monospace label "THE FORMULA" at top with a short navy rule
- The formula rendered large in navy serif, centered, like a print-journal equation
- Below the formula: a thin horizontal rule spanning about 60% of the width
- Below the rule: a routing condition in monospace on a light gray background strip with a red-orange left border
- Below that: a small italic caption in gray

TEXT TO RENDER:
Label: "THE FORMULA"
Formula: "PQL Score = (Production x 0.4) + (Enterprise x 0.4) + (Engagement x 0.2)"
Routing condition: "IF score > 70 AND team_size >= 3  ->  Route to Sales"
Caption: "Start simple. Advanced variants split developer-led vs. top-down."

{TEXT_RULES}"""
    },
    {
        "id": "slide-13-two-routes",
        "prompt": f"""{STYLE}

Create an infographic slide showing two routing paths in a two-column split layout.

COMPOSITION:
- Top: monospace label "ROUTING LOGIC" with a short navy rule, then serif headline "Two Motions"
- At top center: a small diagram of one arrow splitting into two arrows (one going left, one going right) — a diverging path symbol
- Two equal columns separated by a thin vertical rule down the center
- Each column has: a bold navy serif header, a monospace condition in gray, a thin horizontal rule, a bold routing action, and a tactical note in gray

LEFT COLUMN:
Header: "Developer-Led"
Condition: "High production usage + low enterprise signals"
Action: "Route to DevRel."
Tactic: "Nurture with technical content. Sales calls kill this motion."

RIGHT COLUMN:
Header: "Top-Down"
Condition: "Enterprise signals + moderate usage"
Action: "Route to Enterprise AE."
Tactic: "Speed matters. Exec evaluations move fast or die."

{TEXT_RULES}"""
    },
    {
        "id": "slide-14-brand-infra",
        "prompt": f"""{STYLE}

EXCEPTION: This slide uses a DARK BACKGROUND. Background: #1B3A5B deep navy. All text is cream #F4EFE6. Red-orange #C9402F for the rule.

Create a section divider slide for brand infrastructure. Dark navy background.

COMPOSITION:
- Monospace label "SYSTEM 03 — BRAND INFRASTRUCTURE" in muted cream
- Red-orange rule beneath
- Large serif headline with an arrow symbol in it
- Italic subtitle in lighter cream
- Below: ten small monospace tags in rounded pill shapes arranged in two rows of five, with semi-transparent backgrounds
- On the left side: a large serif numeral "1" and on the right side: a large serif numeral "10", with thin radiating lines fanning from the 1 toward the 10

TEXT TO RENDER:
Label: "SYSTEM 03 — BRAND INFRASTRUCTURE"
Headline: "One Technical Piece -> Ten Developer Touchpoints"
Subtitle: "Build once. Distribute across every channel where developers discover tools."
Channel tags: "GitHub" "Dev.to" "Reddit" "Discord" "HackerNews" "YouTube" "Blog" "X/Twitter" "Stack Overflow" "Newsletter"

{TEXT_RULES}"""
    },
    {
        "id": "slide-15-github",
        "prompt": f"""{STYLE}

Create an infographic slide about GitHub as a distribution channel.

COMPOSITION:
- Top: monospace label "DISTRIBUTION ENGINE" with a short navy rule
- Large navy serif headline
- A pull quote in italic serif with a red-orange left border bar
- Five monospace tags in rounded gray pill shapes arranged in a horizontal row
- Below: a short supporting paragraph in gray
- Right side: an abstract illustration of stacked rectangular shapes suggesting code repositories — flat, geometric, layered cards or documents

TEXT TO RENDER:
Label: "DISTRIBUTION ENGINE"
Headline: "GitHub as Primary Distribution"
Pull quote: "Developers trust code over copy. A repo with 500 stars is more credible than a case study."
Tags: "SDKs" "Example repos" "Starter templates" "Migration scripts" "Debugging utilities"
Supporting text: "Every repo is a touchpoint. Every fork is a signal. Every issue and PR is engagement you didn't pay for."

{TEXT_RULES}"""
    },
    {
        "id": "slide-16-three-paths",
        "prompt": f"""{STYLE}

Create an infographic slide showing three scenario paths as horizontal flow sequences.

COMPOSITION:
- Top: monospace label "INTEGRATION MAP" with a short navy rule, then serif headline "Three Paths to Pipeline"
- Three horizontal rows stacked vertically, separated by thin warm-gray horizontal rules
- Each row has: a bold serif label on the left (fixed width), then a sequence of small rounded rectangles connected by thin arrows flowing left to right
- The rounded rectangles contain short text labels
- Row 1 final step should have a navy background with cream text (positive outcome)
- Row 2 final step should have a gray background (neutral outcome)
- Row 3 final step should have a red-orange background with white text (high-value outcome)

TEXT TO RENDER:
Label: "INTEGRATION MAP"
Headline: "Three Paths to Pipeline"
Row 1 label: "Path A — High-Quality PQL"
Row 1 steps: "Signup" -> "Integrate" -> "Scale usage" -> "Team invites" -> "Route to Sales"
Row 2 label: "Path B — Low Signal"
Row 2 steps: "Signup" -> "Sandbox only" -> "No team" -> "Nurture, don't call"
Row 3 label: "Path C — GitHub to Enterprise"
Row 3 steps: "Repo discovery" -> "Team sprawl" -> "Director asks 'what is this?'" -> "Inbound enterprise PQL"

{TEXT_RULES}"""
    },
    {
        "id": "slide-17-metrics",
        "prompt": f"""{STYLE}

Create an infographic slide with a 2x2 measurement framework grid.

COMPOSITION:
- Top: monospace label "METRICS" with a short navy rule, then serif headline "What to Measure"
- Below: a 2x2 grid of four equal quadrants separated by thin cross-shaped dividing lines (warm gray)
- Each quadrant has: a monospace uppercase header in navy, and three bullet-point metrics in sans-serif body text
- The grid should be spacious with generous padding inside each cell
- Clean, organized, balanced layout

TEXT TO RENDER:
Label: "METRICS"
Headline: "What to Measure"
Top-left header: "TOP OF FUNNEL"
Top-left metrics: "Work-email signups" "Production integrations" "GitHub stars and forks"
Top-right header: "MIDDLE OF FUNNEL"
Top-right metrics: "PQL-to-opportunity rate" "Time from PQL to close" "Developer-sourced leads"
Bottom-left header: "REVENUE HEALTH"
Bottom-left metrics: "Expansion revenue" "PLG vs. sales-led CAC" "NRR by channel"
Bottom-right header: "DISTRIBUTION"
Bottom-right metrics: "Repo engagement" "Content multiplier ratio" "Organic share of voice"

{TEXT_RULES}"""
    },
    {
        "id": "slide-18-failures",
        "prompt": f"""{STYLE}

EXCEPTION: This slide uses a DARK BACKGROUND. Background: #1B3A5B deep navy. All text is cream #F4EFE6. Red-orange #C9402F for the rule and warning glyphs.

Create an infographic slide showing four failure modes. Dark navy background.

COMPOSITION:
- Top: monospace label "WATCH FOR" in muted cream, red-orange rule, serif headline "Four Failure Modes" in cream
- Four stacked rows separated by thin semi-transparent rules
- Each row has: a small red-orange circle with an X inside it on the left (warning glyph), a bold cream serif diagnosis, and below it a gray fix starting with an arrow
- Clean, structured, warning-style layout

TEXT TO RENDER:
Label: "WATCH FOR"
Headline: "Four Failure Modes"
Row 1 diagnosis: "Loose friction + aggressive PQL thresholds"
Row 1 fix: "-> Phantom pipeline. Tighten signup, add enterprise signals."
Row 2 diagnosis: "Tight friction + weak brand infrastructure"
Row 2 fix: "-> Growth ceiling. Build the content engine."
Row 3 diagnosis: "Strong usage signals + no conversion wall"
Row 3 fix: "-> Permanently free users. Recalibrate limits and feature gates."
Row 4 diagnosis: "B2B marketing tactics aimed at developers"
Row 4 fix: "-> Credibility death. GitHub first. Marketing second."

{TEXT_RULES}"""
    },
    {
        "id": "slide-19-closing",
        "prompt": f"""{STYLE}

Create a closing slide with a powerful pull quote as the centerpiece.

COMPOSITION:
- Centered layout with maximum whitespace
- Small monospace label "THE FRAMEWORK IN ONE SENTENCE" at top
- A short navy rule beneath
- A large italic serif pull quote centered, with a thick red-orange left border bar (4px)
- The quote should be the visual hero of the slide — large, commanding, elegant
- Below the quote: a thin horizontal rule, then the author/brand line in small gray sans-serif
- Elegant, confident, final-word energy

TEXT TO RENDER:
Label: "THE FRAMEWORK IN ONE SENTENCE"
Pull quote: "The orthodoxy says remove all friction. The physics say apply friction where it does useful work, and remove it everywhere else."
Byline: "Nick Talbert — Strategnik — strategnik.com"

{TEXT_RULES}"""
    },
]


def generate_slide(slide: dict, retries: int = 2) -> Path | None:
    """Generate a single slide image via fal.ai."""
    slide_id = slide["id"]
    output_path = OUTPUT_DIR / f"{slide_id}.png"

    if output_path.exists():
        print(f"  {slide_id}: exists, skipping")
        return output_path

    for attempt in range(retries + 1):
        try:
            print(f"  {slide_id}: generating...")
            resp = requests.post(
                FAL_URL,
                headers={
                    "Authorization": f"Key {FAL_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "prompt": slide["prompt"],
                    "image_size": {"width": IMG_W, "height": IMG_H},
                },
                timeout=120,
            )
            resp.raise_for_status()
            data = resp.json()

            if "images" not in data or not data["images"]:
                raise ValueError(f"No images returned: {list(data.keys())}")

            img_url = data["images"][0]["url"]
            img_resp = requests.get(img_url, timeout=60)
            img_resp.raise_for_status()

            img = Image.open(BytesIO(img_resp.content)).convert("RGB")
            img.save(output_path, "PNG")
            print(f"  {slide_id}: saved")
            return output_path

        except Exception as e:
            print(f"  {slide_id}: attempt {attempt+1} failed — {e}")
            if attempt < retries:
                time.sleep(2)
    return None


def assemble_pptx(slide_images: list[tuple[str, Path]], output_path: Path):
    """Build PPTX with full-bleed images."""
    prs = Presentation()
    # Set 16:9 dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    for slide_id, img_path in slide_images:
        if not img_path or not img_path.exists():
            print(f"  Skipping {slide_id} — no image")
            continue

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(output_path))
    print(f"\nPPTX saved: {output_path}")


def save_prompts(slides: list[dict], output_path: Path):
    """Save prompts for reference / manual refinement."""
    prompts = {s["id"]: s["prompt"] for s in slides}
    output_path.write_text(json.dumps(prompts, indent=2))
    print(f"Prompts saved: {output_path}")


def main():
    print("=== PLG Framework Infographic Deck ===")
    print(f"Model: {MODEL}")
    print(f"Slides: {len(SLIDES)}")
    print(f"Output: {OUTPUT_DIR}\n")

    # Save prompts for reference
    save_prompts(SLIDES, OUTPUT_DIR.parent / "prompts.json")

    # Generate all slide images
    print("Generating slide images...")
    results = []
    for slide in SLIDES:
        img_path = generate_slide(slide)
        results.append((slide["id"], img_path))

    generated = sum(1 for _, p in results if p)
    print(f"\nGenerated {generated}/{len(SLIDES)} slides.\n")

    # Assemble PPTX
    pptx_path = OUTPUT_DIR.parent / "PLG-Framework-Infographic-Deck.pptx"
    print("Assembling PPTX...")
    assemble_pptx(results, pptx_path)

    print("\nDone!")


if __name__ == "__main__":
    main()
